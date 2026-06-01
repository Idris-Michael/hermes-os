"""
Build the Legal Pioneer Pipeline LinkedIn showcase deck (.pptx).

Run:
    py -3.10 scripts/build-legal-pioneer-deck.py
    py -3.10 scripts/build-legal-pioneer-deck.py --video-url https://youtu.be/XXX

Inputs:
    Presentation slide deck assets - legal pioneer/*.png  (captured via Playwright)
    scripts/legal-pioneer-captions.json  (auto-generated on first run; user can edit)

Output:
    deliverables/legal-pioneer-pipeline-deck.pptx
"""
from __future__ import annotations

import argparse
import datetime as dt
import json
import sys
from dataclasses import dataclass
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# ─── Paths ─────────────────────────────────────────────────────────────────────
ROOT = Path(__file__).resolve().parent.parent
ASSETS_DIR = ROOT / "Presentation slide deck assets - legal pioneer"
CAPTIONS_FILE = ROOT / "scripts" / "legal-pioneer-captions.json"
OUTPUT = ROOT / "deliverables" / "legal-pioneer-pipeline-deck.pptx"

# ─── Design system (same palette as Hermes deck) ───────────────────────────────
def hex_rgb(s: str) -> RGBColor:
    s = s.lstrip("#")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))

BG       = hex_rgb("#060610")
SURFACE  = hex_rgb("#0E0E1A")
ACCENT   = hex_rgb("#D4A017")
TEXT     = hex_rgb("#E2E8F0")
TEXT_DIM = hex_rgb("#94A3B8")
RULE     = hex_rgb("#1F1F2E")
GREEN    = hex_rgb("#4ADE80")
AMBER    = hex_rgb("#FBBF24")
RED      = hex_rgb("#F87171")
VIOLET   = hex_rgb("#C4B5FD")

FONT_DISPLAY = "Inter SemiBold"
FONT_BODY    = "Inter"
FONT_MONO    = "JetBrains Mono"

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# ─── Project copy ──────────────────────────────────────────────────────────────
PROJECT_TITLE     = "Legal Pioneer Pipeline"
PROJECT_SUBTITLE  = "AI-grounded contract review for in-house legal teams"
PROJECT_CONTEXT   = "PortSwigger AI Pioneer Submission · May 2026 · Idris-Michael Bakare"

THESIS = (
    "Small legal teams don't drown in contract complexity.\n"
    "They drown in volume — the end-of-quarter spike that turns 10 NDAs a week into 35."
)
THESIS_SUB = (
    "The Legal Pioneer Pipeline handles the 80% of contracts that are boilerplate so "
    "the legal team can focus on the 20% that genuinely needs them — without ever guessing."
)

METRICS = [
    ("62%",         "Lawyer time saved on routine NDAs and order forms",                    "Time reclaimed"),
    ("< £0.01",     "Cost per contract review at 10 contracts/day (Claude + Gemini APIs)",  "Per-contract cost"),
    ("100%",        "Auditable decisions — every classification traces to a brain.md rule", "Explainability"),
]

PILLARS = [
    ("Grounded, not guessing",
     "Every clause matched to a specific policy in brain.md. UNMAPPED clauses are flagged for "
     "humans — never silently classified."),
    ("Legal team owns the AI",
     "brain.md is plain Markdown. The legal team edits the playbook themselves. No engineering "
     "dependency, no ticket queue."),
    ("Two-agent verification",
     "Architect (Claude) classifies and proposes rewrites. Auditor (Gemini) runs in parallel and "
     "catches misclassifications before output."),
    ("The EOQ override",
     "When the queue hits crisis, the legal lead adds one Markdown block to brain.md. Pipeline "
     "picks it up on the next run. No deploy, no ticket."),
]

TIER_TABLE = [
    ("Green",    "Standard boilerplate, no material deviation",   "Auto-approve with log",                "0 min"),
    ("Amber",    "Minor variation, known edit exists",             "Propose rewrite, await approval",      "2-3 min"),
    ("Red",      "Novel clause, material risk, no precedent",      "Escalate immediately with context",    "Full review"),
    ("UNMAPPED", "No matching policy in brain.md",                 "Flag for human — never auto-decide",   "Full review"),
]

STACK_CHIPS = [
    "Claude Sonnet 4.6", "Gemini 1.5 Pro", "NotebookLM", "Obsidian",
    "Python 3.10", "PyPDF2", "python-docx", "Anthropic SDK", "Google Gen AI SDK",
    "Markdown Playbook", "Agent Orchestration", "Multi-Agent Verification",
    "Hallucination Guardrails", "Explainable AI", "Audit Logging",
]

CONTACT = {
    "name":     "Idris-Michael Bakare",
    "role":     "AI Pioneer · Severus Connects · London",
    "email":    "hello@severusconnects.com",
    "linkedin": "linkedin.com/in/idris-michael-bakare",
    "cta":      "Available to walk through the pipeline live.",
}

# ─── Helpers ───────────────────────────────────────────────────────────────────
@dataclass(frozen=True)
class Asset:
    path: Path
    width: int
    height: int

def discover_assets() -> list[Asset]:
    if not ASSETS_DIR.exists():
        sys.exit(f"Assets folder not found: {ASSETS_DIR}\nRun capture-legal-pioneer-assets.py first.")
    pngs = sorted(ASSETS_DIR.glob("*.png"), key=lambda p: p.name)
    if not pngs:
        sys.exit(f"No PNGs found in: {ASSETS_DIR}")
    out: list[Asset] = []
    for p in pngs:
        with Image.open(p) as im:
            out.append(Asset(path=p, width=im.width, height=im.height))
    return out

DEFAULT_CAPTIONS: dict[str, str] = {
    "01-demo-full-report.png":            "Full pipeline output — one auto-generated review of a real-world mutual NDA. Green, Amber, Red, and UNMAPPED tiers visible at a glance.",
    "02-demo-header-and-escalation.png":  "Above the fold — pipeline metadata, escalation banner, and the 3/2/1/1 clause distribution. The lawyer's 30-second triage.",
    "03-clause-01.png":                   "Green tier — Confidentiality Definition matches the standard position exactly. Auto-approved against brain.md § 1.",
    "03-clause-02.png":                   "Green tier — Term and Termination exceeds the standard position (strictly more protective). Auto-approved.",
    "03-clause-03.png":                   "Green tier — Return / Destruction matches the playbook. Auto-approved with full policy traceability.",
    "03-clause-04.png":                   "Amber tier — Permitted Disclosure proposes 'all group companies' without limitation. Pipeline proposes the standard-position rewrite.",
    "03-clause-05.png":                   "Amber tier — Limitation of Liability cap below threshold. Rewrite ready for one-click lawyer approval.",
    "03-clause-06.png":                   "Red tier — Governing Law proposes Delaware. Hard escalation: no policy permits non-UK jurisdiction.",
    "03-clause-07.png":                   "UNMAPPED tier — Residual Knowledge clause has no policy in brain.md. Flagged for human, never silently classified.",
    "10-source-brain.png":                "brain.md — the AI's grounding source. Extracted from 50 historical NDAs via NotebookLM, edited by the legal team in plain Markdown.",
    "11-source-agents.png":               "agents.md — Architect (Claude) classifies and rewrites. Auditor (Gemini) verifies. Orchestrator runs both in parallel.",
    "12-source-pipeline.png":             "review_pipeline.py — CLI orchestration. Two API keys, one PDF or DOCX, one HTML report. Runs in seconds.",
}

def load_or_seed_captions(assets: list[Asset]) -> dict[str, str]:
    if CAPTIONS_FILE.exists():
        try:
            existing: dict[str, str] = json.loads(CAPTIONS_FILE.read_text(encoding="utf-8"))
        except json.JSONDecodeError:
            existing = {}
    else:
        existing = {}
    seeded = {a.path.name: existing.get(a.path.name) or DEFAULT_CAPTIONS.get(a.path.name, "")
              for a in assets}
    CAPTIONS_FILE.parent.mkdir(parents=True, exist_ok=True)
    CAPTIONS_FILE.write_text(json.dumps(seeded, indent=2, sort_keys=True) + "\n", encoding="utf-8")
    return seeded

def fill_solid(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_text(slide, left_in: float, top_in: float, width_in: float, height_in: float,
             text: str, *, font: str = FONT_BODY, size_pt: int = 18, color: RGBColor = TEXT,
             bold: bool = False, align: int = PP_ALIGN.LEFT, anchor: int = MSO_ANCHOR.TOP) -> None:
    tb = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size_pt)
    f.bold = bold
    f.color.rgb = color

def paint_bg(slide) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SLIDE_W_IN), Inches(SLIDE_H_IN))
    fill_solid(bg, BG)

def add_accent_bar(slide, left_in: float = 0.7, top_in: float = 0.7, width_in: float = 0.6) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left_in), Inches(top_in), Inches(width_in), Inches(0.08))
    fill_solid(bar, ACCENT)

def add_page_footer(slide, label: str) -> None:
    add_text(slide, 0.5, SLIDE_H_IN - 0.45, SLIDE_W_IN - 1.0, 0.3,
             label, font=FONT_MONO, size_pt=9, color=TEXT_DIM)

# ─── Slide builders ────────────────────────────────────────────────────────────
def slide_cover(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    add_accent_bar(s, 0.7, 2.2, 1.6)
    add_text(s, 0.7, 2.5, 12, 1.4, PROJECT_TITLE,
             font=FONT_DISPLAY, size_pt=72, color=TEXT, bold=True)
    add_text(s, 0.7, 3.95, 12, 0.7, PROJECT_SUBTITLE,
             font=FONT_DISPLAY, size_pt=24, color=ACCENT)
    add_text(s, 0.7, 4.8, 12, 0.5,
             "A working prototype of contract review that doesn't guess.",
             font=FONT_BODY, size_pt=16, color=TEXT_DIM)
    add_text(s, 0.7, 6.6, 12, 0.4, PROJECT_CONTEXT,
             font=FONT_MONO, size_pt=11, color=TEXT_DIM)

def slide_thesis(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    add_accent_bar(s)
    add_text(s, 0.7, 1.0, 12, 0.6, "01 · The problem",
             font=FONT_MONO, size_pt=12, color=ACCENT)
    add_text(s, 0.7, 2.0, 12, 3.0, THESIS,
             font=FONT_DISPLAY, size_pt=40, color=TEXT, bold=True)
    add_text(s, 0.7, 5.4, 12, 1.4, THESIS_SUB,
             font=FONT_BODY, size_pt=18, color=TEXT_DIM)
    add_page_footer(s, "legal-pioneer · thesis")

def slide_pillars(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    add_accent_bar(s)
    add_text(s, 0.7, 1.0, 12, 0.6, "02 · How it works",
             font=FONT_MONO, size_pt=12, color=ACCENT)
    add_text(s, 0.7, 1.5, 12, 0.7, "Four design decisions.",
             font=FONT_DISPLAY, size_pt=32, color=TEXT, bold=True)
    # 2 x 2 grid
    coords = [(0.7, 2.6), (6.9, 2.6), (0.7, 4.95), (6.9, 4.95)]
    for (title, body), (x, y) in zip(PILLARS, coords):
        card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(5.7), Inches(2.1))
        card.fill.solid()
        card.fill.fore_color.rgb = SURFACE
        card.line.color.rgb = RULE
        card.line.width = Pt(0.75)
        add_text(s, x + 0.3, y + 0.25, 5.3, 0.5, title,
                 font=FONT_DISPLAY, size_pt=18, color=ACCENT, bold=True)
        add_text(s, x + 0.3, y + 0.85, 5.3, 1.2, body,
                 font=FONT_BODY, size_pt=13, color=TEXT)
    add_page_footer(s, "legal-pioneer · pillars")

def slide_tier_table(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    add_accent_bar(s)
    add_text(s, 0.7, 1.0, 12, 0.6, "03 · The triage system",
             font=FONT_MONO, size_pt=12, color=ACCENT)
    add_text(s, 0.7, 1.5, 12, 0.7, "Every clause lands in one of four tiers.",
             font=FONT_DISPLAY, size_pt=28, color=TEXT, bold=True)

    headers = ["Tier", "Description", "Pipeline action", "Lawyer time"]
    col_x = [0.7, 2.2, 6.9, 11.2]
    col_w = [1.5, 4.7, 4.3, 1.7]
    header_top = 3.0
    add_text(s, col_x[0], header_top, col_w[0], 0.4, headers[0],
             font=FONT_MONO, size_pt=11, color=TEXT_DIM)
    add_text(s, col_x[1], header_top, col_w[1], 0.4, headers[1],
             font=FONT_MONO, size_pt=11, color=TEXT_DIM)
    add_text(s, col_x[2], header_top, col_w[2], 0.4, headers[2],
             font=FONT_MONO, size_pt=11, color=TEXT_DIM)
    add_text(s, col_x[3], header_top, col_w[3], 0.4, headers[3],
             font=FONT_MONO, size_pt=11, color=TEXT_DIM)
    rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(header_top + 0.45),
                               Inches(SLIDE_W_IN - 1.4), Inches(0.02))
    fill_solid(rule, RULE)

    tier_colors = {"Green": GREEN, "Amber": AMBER, "Red": RED, "UNMAPPED": VIOLET}
    row_h = 0.85
    for i, row in enumerate(TIER_TABLE):
        y = header_top + 0.7 + i * row_h
        tier_color = tier_colors[row[0]]
        # tier name
        add_text(s, col_x[0], y, col_w[0], 0.5, row[0],
                 font=FONT_DISPLAY, size_pt=16, color=tier_color, bold=True)
        add_text(s, col_x[1], y + 0.05, col_w[1], 0.7, row[1],
                 font=FONT_BODY, size_pt=13, color=TEXT)
        add_text(s, col_x[2], y + 0.05, col_w[2], 0.7, row[2],
                 font=FONT_BODY, size_pt=13, color=TEXT)
        add_text(s, col_x[3], y + 0.05, col_w[3], 0.7, row[3],
                 font=FONT_MONO, size_pt=13, color=TEXT_DIM)
    add_page_footer(s, "legal-pioneer · triage")

def slide_metric(prs: Presentation, big: str, sub: str, label: str) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    add_accent_bar(s)
    add_text(s, 0.7, 1.0, 12, 0.6, label,
             font=FONT_MONO, size_pt=12, color=ACCENT)
    add_text(s, 0.7, 2.4, 12, 2.5, big,
             font=FONT_DISPLAY, size_pt=130, color=ACCENT, bold=True)
    add_text(s, 0.7, 5.5, 12, 1.2, sub,
             font=FONT_BODY, size_pt=22, color=TEXT)
    add_page_footer(s, "legal-pioneer · impact")

def slide_screenshot(prs: Presentation, index: int, total: int, asset: Asset, caption: str) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    img_area_top = 0.4
    img_area_height = SLIDE_H_IN - img_area_top - 0.9
    img_area_width = SLIDE_W_IN - 1.0
    img_area_left = 0.5

    aspect = asset.width / asset.height
    area_aspect = img_area_width / img_area_height
    if aspect >= area_aspect:
        new_w = img_area_width
        new_h = new_w / aspect
    else:
        new_h = img_area_height
        new_w = new_h * aspect
    left = img_area_left + (img_area_width - new_w) / 2
    top = img_area_top + (img_area_height - new_h) / 2
    s.shapes.add_picture(str(asset.path), Inches(left), Inches(top), Inches(new_w), Inches(new_h))

    strip_top = SLIDE_H_IN - 0.85
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.5), Inches(strip_top),
                              Inches(0.08), Inches(0.45))
    fill_solid(bar, ACCENT)
    add_text(s, 0.75, strip_top - 0.05, 11.5, 0.55,
             caption if caption else "—",
             font=FONT_DISPLAY, size_pt=14,
             color=TEXT if caption else TEXT_DIM)
    add_text(s, SLIDE_W_IN - 1.2, strip_top + 0.05, 0.8, 0.3,
             f"{index:02d}/{total:02d}",
             font=FONT_MONO, size_pt=10, color=TEXT_DIM, align=PP_ALIGN.RIGHT)

def slide_eoq_override(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    add_accent_bar(s)
    add_text(s, 0.7, 1.0, 12, 0.6, "04 · The end-of-quarter override",
             font=FONT_MONO, size_pt=12, color=ACCENT)
    add_text(s, 0.7, 1.5, 12, 1.0,
             "When the queue hits crisis, the legal team\nedits one Markdown block.",
             font=FONT_DISPLAY, size_pt=28, color=TEXT, bold=True)

    code = ("## EOQ Override — active until 2026-06-30\n"
            "Auto-approve all mutual NDAs under £10,000 where\n"
            "counterparty is a UK-registered direct customer.\n\n"
            "Authorised by: [Legal Lead Name]")
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(3.6),
                              Inches(SLIDE_W_IN - 1.4), Inches(2.4))
    box.fill.solid()
    box.fill.fore_color.rgb = SURFACE
    box.line.color.rgb = ACCENT
    box.line.width = Pt(1.0)
    add_text(s, 1.0, 3.85, SLIDE_W_IN - 2.0, 2.1, code,
             font=FONT_MONO, size_pt=15, color=GREEN)

    add_text(s, 0.7, 6.4, 12, 0.5,
             "The pipeline picks it up on the next run. No deployment. No ticket. No waiting.",
             font=FONT_BODY, size_pt=16, color=TEXT_DIM)
    add_page_footer(s, "legal-pioneer · eoq-override")

def slide_tech_chips(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    add_accent_bar(s)
    add_text(s, 0.7, 1.0, 12, 0.6, "Stack proof",
             font=FONT_MONO, size_pt=12, color=ACCENT)
    add_text(s, 0.7, 1.5, 12, 0.8, "What it's built with.",
             font=FONT_DISPLAY, size_pt=36, color=TEXT, bold=True)

    chip_h = 0.5
    pad_x = 0.2
    pad_y = 0.2
    char_w = 0.105
    left = 0.7
    top = 2.8
    cur_x = left
    cur_y = top
    max_x = SLIDE_W_IN - 0.7

    for chip in STACK_CHIPS:
        w = max(1.1, len(chip) * char_w + 0.5)
        if cur_x + w > max_x:
            cur_x = left
            cur_y += chip_h + pad_y
        rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(cur_x), Inches(cur_y), Inches(w), Inches(chip_h))
        rect.fill.solid()
        rect.fill.fore_color.rgb = SURFACE
        rect.line.color.rgb = ACCENT
        rect.line.width = Pt(0.75)
        tf = rect.text_frame
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = chip
        run.font.name = FONT_DISPLAY
        run.font.size = Pt(13)
        run.font.color.rgb = ACCENT
        cur_x += w + pad_x
    add_page_footer(s, "legal-pioneer · stack")

def slide_cta(prs: Presentation, video_url: str | None) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    add_accent_bar(s)
    add_text(s, 0.7, 1.0, 12, 0.6, "Watch the 60-second walkthrough",
             font=FONT_MONO, size_pt=12, color=ACCENT)
    add_text(s, 0.7, 2.0, 12, 1.5, "See it review a contract.",
             font=FONT_DISPLAY, size_pt=52, color=TEXT, bold=True)
    add_text(s, 0.7, 3.6, 12, 1,
             "Contract goes in. Tiered report comes out — with every decision\n"
             "traceable to a policy line in brain.md.",
             font=FONT_BODY, size_pt=18, color=TEXT_DIM)
    url = video_url or "YOUTUBE_URL_PENDING"
    add_text(s, 0.7, 5.3, 12, 0.5, "Link:",
             font=FONT_MONO, size_pt=12, color=TEXT_DIM)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(5.7), Inches(12), Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = url
    run.font.name = FONT_MONO
    run.font.size = Pt(22)
    run.font.color.rgb = ACCENT
    run.font.bold = True
    if video_url:
        run.hyperlink.address = video_url
    add_page_footer(s, "legal-pioneer · cta")

def slide_contact(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    add_accent_bar(s)
    add_text(s, 0.7, 1.0, 12, 0.6, "Contact",
             font=FONT_MONO, size_pt=12, color=ACCENT)
    add_text(s, 0.7, 2.2, 12, 1.0, CONTACT["name"],
             font=FONT_DISPLAY, size_pt=54, color=TEXT, bold=True)
    add_text(s, 0.7, 3.4, 12, 0.7, CONTACT["role"],
             font=FONT_BODY, size_pt=22, color=TEXT_DIM)
    add_text(s, 0.7, 5.0, 12, 0.5, CONTACT["email"],
             font=FONT_MONO, size_pt=18, color=ACCENT)
    add_text(s, 0.7, 5.6, 12, 0.5, CONTACT["linkedin"],
             font=FONT_MONO, size_pt=18, color=ACCENT)
    add_text(s, 0.7, 6.6, 12, 0.4, CONTACT["cta"],
             font=FONT_DISPLAY, size_pt=18, color=TEXT)

# ─── Build ─────────────────────────────────────────────────────────────────────
def build(video_url: str | None) -> None:
    assets = discover_assets()
    captions = load_or_seed_captions(assets)

    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    # Framing
    slide_cover(prs)
    slide_thesis(prs)
    slide_pillars(prs)
    slide_tier_table(prs)

    # Metric 1 (time saved)
    slide_metric(prs, METRICS[0][0], METRICS[0][1], METRICS[0][2])

    # Visual proof — walkthrough screenshots
    total = len(assets)
    for i, a in enumerate(assets, start=1):
        slide_screenshot(prs, i, total, a, captions.get(a.path.name, ""))

    # Metric 2 (cost)
    slide_metric(prs, METRICS[1][0], METRICS[1][1], METRICS[1][2])

    # EOQ override slide
    slide_eoq_override(prs)

    # Metric 3 (explainability)
    slide_metric(prs, METRICS[2][0], METRICS[2][1], METRICS[2][2])

    # Closers
    slide_tech_chips(prs)
    slide_cta(prs, video_url)
    slide_contact(prs)

    cp = prs.core_properties
    cp.author = "Idris-Michael Bakare"
    cp.title = "Legal Pioneer Pipeline — Portfolio Deck"
    cp.modified = dt.datetime(2026, 1, 1, 0, 0, 0)
    cp.created  = dt.datetime(2026, 1, 1, 0, 0, 0)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"OK -> {OUTPUT}")
    print(f"   {len(prs.slides)} slides | {len(assets)} screenshots | captions: {CAPTIONS_FILE}")
    if not video_url:
        print("   video URL placeholder -> rerun with --video-url <url>")

def main() -> None:
    ap = argparse.ArgumentParser(description="Build the Legal Pioneer Pipeline showcase deck.")
    ap.add_argument("--video-url", default=None, help="YouTube/Loom URL for the CTA slide.")
    args = ap.parse_args()
    build(args.video_url)

if __name__ == "__main__":
    main()
