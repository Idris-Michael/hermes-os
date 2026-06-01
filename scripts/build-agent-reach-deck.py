"""
Build the Agent-Reach LinkedIn showcase deck (.pptx).

Run:
    py -3.10 scripts/build-agent-reach-deck.py
    py -3.10 scripts/build-agent-reach-deck.py --video-url https://youtu.be/XXX

Inputs:
    Presentation slide deck assets - agent-reach/*.png
    scripts/agent-reach-captions.json (auto-generated on first run; user can edit)

Output:
    deliverables/agent-reach-deck.pptx
"""
from __future__ import annotations

import argparse
import datetime as dt
import json
import sys
from dataclasses import dataclass
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# Paths
ROOT = Path(__file__).resolve().parent.parent
ASSETS_DIR = ROOT / "Presentation slide deck assets - agent-reach"
CAPTIONS_FILE = ROOT / "scripts" / "agent-reach-captions.json"
OUTPUT = ROOT / "deliverables" / "agent-reach-deck.pptx"

# Design system -- identical palette to Legal Pioneer
def hex_rgb(s: str) -> RGBColor:
    s = s.lstrip("#")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))

BG       = hex_rgb("#060610")
SURFACE  = hex_rgb("#0E0E1A")
ACCENT   = hex_rgb("#D4A017")
TEXT     = hex_rgb("#E2E8F0")
TEXT_DIM = hex_rgb("#94A3B8")
RULE     = hex_rgb("#1F1F2E")

FONT_DISPLAY = "Inter SemiBold"
FONT_BODY    = "Inter"
FONT_MONO    = "JetBrains Mono"

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# Project copy
PROJECT_TITLE     = "Agent-Reach"
PROJECT_SUBTITLE  = "Web Toolkit for AI Agents"
PROJECT_CONTEXT   = "Open-source Python package . 2026 . Idris-Michael Bakare"

THESIS = (
    "Most AI coding agents have rich filesystem access\n"
    "and almost zero access to the open web."
)
THESIS_SUB = (
    "Agent-Reach is the smallest possible Python package that closes that gap. Drop it into any "
    "agent runtime and the agent can read YouTube transcripts, scrape Reddit threads, parse GitHub "
    "repos, and tail RSS feeds -- without a paid API key per platform."
)

METRICS = [
    ("16",      "Platforms supported from a single pip install -- web, video, social, code, news",  "Coverage"),
    ("7",       "Zero-config channels that work the moment the package is installed",               "Out-of-the-box"),
    ("2",       "Shipping languages -- the README mirrors Mandarin and English for day-one reach",  "Internationalised"),
]

PILLARS = [
    ("Drop-in skill",
     "One pip install registers a SKILL.md the agent reads automatically. No bespoke prompt "
     "engineering -- the agent just knows it can now read Twitter, YouTube, Reddit, and more."),
    ("Zero per-platform setup",
     "Seven of the sixteen channels work the second install finishes. The rest configure with one "
     "sentence to the agent: 'log me into GitHub', 'configure Twitter'."),
    ("Built for agent runtimes",
     "Compatible with Claude Code, OpenClaw, Cursor, Windsurf, Codex -- anything that can run a "
     "shell command. Backends are real CLIs (yt-dlp, gh, bird) not bespoke wrappers."),
    ("Open source and pluggable",
     "Every platform is one file in channels/. Don't trust a backend? Swap it. The agent_reach "
     "doctor command tells the operator which channels are healthy and which need work."),
]

# The sixteen platforms from the README support matrix.
STACK_CHIPS = [
    "Web", "YouTube", "RSS", "Whole-Web Search",
    "GitHub", "Twitter/X", "Bilibili", "Reddit",
    "Xiaohongshu", "Douyin", "LinkedIn",
    "WeChat MP", "Weibo", "V2EX", "Xueqiu",
    "Xiaoyuzhou Podcasts",
]

CONTACT = {
    "name":     "Idris-Michael Bakare",
    "role":     "Severus Connects . London",
    "email":    "hello@severusconnects.com",
    "linkedin": "linkedin.com/in/idris-michael-bakare",
    "cta":      "Available to walk through Agent-Reach live.",
}

# Helpers
@dataclass(frozen=True)
class Asset:
    path: Path
    width: int
    height: int

def discover_assets() -> list[Asset]:
    if not ASSETS_DIR.exists():
        sys.exit(f"Assets folder not found: {ASSETS_DIR}\nRun capture-agent-reach-assets.py first.")
    pngs = sorted(ASSETS_DIR.glob("*.png"), key=lambda p: p.name)
    if not pngs:
        sys.exit(f"No PNGs found in: {ASSETS_DIR}")
    out: list[Asset] = []
    for p in pngs:
        with Image.open(p) as im:
            out.append(Asset(path=p, width=im.width, height=im.height))
    return out

DEFAULT_CAPTIONS: dict[str, str] = {
    "01-pitch-readme.png":       "The README -- one install, sixteen platforms. Mandarin original; English mirror in docs/.",
    "02-package-pyproject.png":  "pyproject.toml -- published to PyPI as agent-reach. Minimal core deps, optional MCP and browser extras.",
    "03-api-surface.png":        "agent_reach/__init__.py -- one public class, AgentReach. Everything else is pluggable channels.",
    "04-support-matrix.png":     "The support matrix -- seven zero-config channels, the rest configure with one sentence to the agent.",
    "05-diagnostic-doctor.png":  "agent-reach doctor -- one command reports which channels are healthy and which need attention.",
    "06-cli-entrypoint.png":     "agent_reach/cli.py -- install, doctor, configure, uninstall. The whole operator surface.",
    "07-channels-registry.png":  "channels/__init__.py -- each platform is one file. Swap a backend by editing one file.",
    "08-channel-twitter.png":    "channels/twitter.py -- a channel only declares its backend (bird CLI) and how to health-check it.",
    "09-i18n-readme-zh.png":     "The Mandarin README -- shipped from day one alongside the English mirror.",
    "10-logo.png":               "Agent-Reach mark -- the open-source identity used on GitHub and PyPI.",
}

def load_or_seed_captions(assets: list[Asset]) -> dict[str, str]:
    if CAPTIONS_FILE.exists():
        try:
            existing: dict[str, str] = json.loads(CAPTIONS_FILE.read_text(encoding="utf-8"))
        except json.JSONDecodeError:
            existing = {}
    else:
        existing = {}
    seeded = {a.path.name: existing.get(a.path.name) or DEFAULT_CAPTIONS.get(a.path.name, "")
              for a in assets}
    CAPTIONS_FILE.parent.mkdir(parents=True, exist_ok=True)
    CAPTIONS_FILE.write_text(json.dumps(seeded, indent=2, sort_keys=True) + "\n", encoding="utf-8")
    return seeded

def fill_solid(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_text(slide, left_in: float, top_in: float, width_in: float, height_in: float,
             text: str, *, font: str = FONT_BODY, size_pt: int = 18, color: RGBColor = TEXT,
             bold: bool = False, align: int = PP_ALIGN.LEFT, anchor: int = MSO_ANCHOR.TOP) -> None:
    tb = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size_pt)
    f.bold = bold
    f.color.rgb = color

def paint_bg(slide) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SLIDE_W_IN), Inches(SLIDE_H_IN))
    fill_solid(bg, BG)

def add_accent_bar(slide, left_in: float = 0.7, top_in: float = 0.7, width_in: float = 0.6) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left_in), Inches(top_in),
                                  Inches(width_in), Inches(0.08))
    fill_solid(bar, ACCENT)

def add_page_footer(slide, label: str) -> None:
    add_text(slide, 0.5, SLIDE_H_IN - 0.45, SLIDE_W_IN - 1.0, 0.3,
             label, font=FONT_MONO, size_pt=9, color=TEXT_DIM)

# Slide builders
def slide_cover(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    add_accent_bar(s, 0.7, 2.2, 1.6)
    add_text(s, 0.7, 2.5, 12, 1.4, PROJECT_TITLE,
             font=FONT_DISPLAY, size_pt=84, color=TEXT, bold=True)
    add_text(s, 0.7, 4.05, 12, 0.7, PROJECT_SUBTITLE,
             font=FONT_DISPLAY, size_pt=26, color=ACCENT)
    add_text(s, 0.7, 4.9, 12, 0.5,
             "One install. Sixteen platforms. Give Claude, OpenClaw, or Windsurf agents the open web.",
             font=FONT_BODY, size_pt=16, color=TEXT_DIM)
    add_text(s, 0.7, 6.6, 12, 0.4, PROJECT_CONTEXT,
             font=FONT_MONO, size_pt=11, color=TEXT_DIM)

def slide_thesis(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    add_accent_bar(s)
    add_text(s, 0.7, 1.0, 12, 0.6, "01 . The gap",
             font=FONT_MONO, size_pt=12, color=ACCENT)
    add_text(s, 0.7, 2.0, 12, 3.0, THESIS,
             font=FONT_DISPLAY, size_pt=42, color=TEXT, bold=True)
    add_text(s, 0.7, 5.4, 12, 1.4, THESIS_SUB,
             font=FONT_BODY, size_pt=18, color=TEXT_DIM)
    add_page_footer(s, "agent-reach . thesis")

def slide_pillars(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    add_accent_bar(s)
    add_text(s, 0.7, 1.0, 12, 0.6, "02 . How it works",
             font=FONT_MONO, size_pt=12, color=ACCENT)
    add_text(s, 0.7, 1.5, 12, 0.7, "Four design decisions.",
             font=FONT_DISPLAY, size_pt=32, color=TEXT, bold=True)
    coords = [(0.7, 2.6), (6.9, 2.6), (0.7, 4.95), (6.9, 4.95)]
    for (title, body), (x, y) in zip(PILLARS, coords):
        card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(5.7), Inches(2.1))
        card.fill.solid()
        card.fill.fore_color.rgb = SURFACE
        card.line.color.rgb = RULE
        card.line.width = Pt(0.75)
        add_text(s, x + 0.3, y + 0.25, 5.3, 0.5, title,
                 font=FONT_DISPLAY, size_pt=18, color=ACCENT, bold=True)
        add_text(s, x + 0.3, y + 0.85, 5.3, 1.2, body,
                 font=FONT_BODY, size_pt=13, color=TEXT)
    add_page_footer(s, "agent-reach . pillars")

def slide_metric(prs: Presentation, big: str, sub: str, label: str) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    add_accent_bar(s)
    add_text(s, 0.7, 1.0, 12, 0.6, label,
             font=FONT_MONO, size_pt=12, color=ACCENT)
    add_text(s, 0.7, 2.4, 12, 2.5, big,
             font=FONT_DISPLAY, size_pt=180, color=ACCENT, bold=True)
    add_text(s, 0.7, 5.5, 12, 1.2, sub,
             font=FONT_BODY, size_pt=22, color=TEXT)
    add_page_footer(s, "agent-reach . impact")

def slide_screenshot(prs: Presentation, index: int, total: int, asset: Asset, caption: str) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    img_area_top = 0.4
    img_area_height = SLIDE_H_IN - img_area_top - 0.9
    img_area_width = SLIDE_W_IN - 1.0
    img_area_left = 0.5

    aspect = asset.width / asset.height
    area_aspect = img_area_width / img_area_height
    if aspect >= area_aspect:
        new_w = img_area_width
        new_h = new_w / aspect
    else:
        new_h = img_area_height
        new_w = new_h * aspect
    left = img_area_left + (img_area_width - new_w) / 2
    top = img_area_top + (img_area_height - new_h) / 2
    s.shapes.add_picture(str(asset.path), Inches(left), Inches(top), Inches(new_w), Inches(new_h))

    strip_top = SLIDE_H_IN - 0.85
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.5), Inches(strip_top),
                              Inches(0.08), Inches(0.45))
    fill_solid(bar, ACCENT)
    add_text(s, 0.75, strip_top - 0.05, 11.5, 0.55,
             caption if caption else "--",
             font=FONT_DISPLAY, size_pt=14,
             color=TEXT if caption else TEXT_DIM)
    add_text(s, SLIDE_W_IN - 1.2, strip_top + 0.05, 0.8, 0.3,
             f"{index:02d}/{total:02d}",
             font=FONT_MONO, size_pt=10, color=TEXT_DIM, align=PP_ALIGN.RIGHT)

def slide_tech_chips(prs: Presentation, title: str, eyebrow: str) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    add_accent_bar(s)
    add_text(s, 0.7, 1.0, 12, 0.6, eyebrow,
             font=FONT_MONO, size_pt=12, color=ACCENT)
    add_text(s, 0.7, 1.5, 12, 0.8, title,
             font=FONT_DISPLAY, size_pt=36, color=TEXT, bold=True)

    chip_h = 0.55
    pad_x = 0.2
    pad_y = 0.22
    char_w = 0.115
    left = 0.7
    top = 2.8
    cur_x = left
    cur_y = top
    max_x = SLIDE_W_IN - 0.7

    for chip in STACK_CHIPS:
        w = max(1.2, len(chip) * char_w + 0.6)
        if cur_x + w > max_x:
            cur_x = left
            cur_y += chip_h + pad_y
        rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(cur_x), Inches(cur_y), Inches(w), Inches(chip_h))
        rect.fill.solid()
        rect.fill.fore_color.rgb = SURFACE
        rect.line.color.rgb = ACCENT
        rect.line.width = Pt(0.75)
        tf = rect.text_frame
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = chip
        run.font.name = FONT_DISPLAY
        run.font.size = Pt(13)
        run.font.color.rgb = ACCENT
        cur_x += w + pad_x
    add_page_footer(s, "agent-reach . platforms")

def slide_cta(prs: Presentation, video_url: str | None) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    add_accent_bar(s)
    add_text(s, 0.7, 1.0, 12, 0.6, "Try it",
             font=FONT_MONO, size_pt=12, color=ACCENT)
    add_text(s, 0.7, 2.0, 12, 1.5, "pip install agent-reach",
             font=FONT_MONO, size_pt=46, color=TEXT, bold=True)
    add_text(s, 0.7, 3.6, 12, 1,
             "Then point your agent at the install doc and let it configure itself.\n"
             "agent-reach doctor will tell you exactly what works.",
             font=FONT_BODY, size_pt=18, color=TEXT_DIM)
    url = video_url or "github.com/Panniantong/agent-reach"
    add_text(s, 0.7, 5.3, 12, 0.5, "Link:",
             font=FONT_MONO, size_pt=12, color=TEXT_DIM)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(5.7), Inches(12), Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = url
    run.font.name = FONT_MONO
    run.font.size = Pt(22)
    run.font.color.rgb = ACCENT
    run.font.bold = True
    if video_url:
        run.hyperlink.address = video_url
    add_page_footer(s, "agent-reach . cta")

def slide_contact(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    add_accent_bar(s)
    add_text(s, 0.7, 1.0, 12, 0.6, "Contact",
             font=FONT_MONO, size_pt=12, color=ACCENT)
    add_text(s, 0.7, 2.2, 12, 1.0, CONTACT["name"],
             font=FONT_DISPLAY, size_pt=54, color=TEXT, bold=True)
    add_text(s, 0.7, 3.4, 12, 0.7, CONTACT["role"],
             font=FONT_BODY, size_pt=22, color=TEXT_DIM)
    add_text(s, 0.7, 5.0, 12, 0.5, CONTACT["email"],
             font=FONT_MONO, size_pt=18, color=ACCENT)
    add_text(s, 0.7, 5.6, 12, 0.5, CONTACT["linkedin"],
             font=FONT_MONO, size_pt=18, color=ACCENT)
    add_text(s, 0.7, 6.6, 12, 0.4, CONTACT["cta"],
             font=FONT_DISPLAY, size_pt=18, color=TEXT)

# Build
def build(video_url: str | None) -> None:
    assets = discover_assets()
    captions = load_or_seed_captions(assets)

    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    # Framing
    slide_cover(prs)
    slide_thesis(prs)
    slide_pillars(prs)
    slide_tech_chips(prs, "What's in the box.", "03 . Sixteen platforms")

    # Metric 1 (coverage)
    slide_metric(prs, METRICS[0][0], METRICS[0][1], METRICS[0][2])

    # Visual proof -- source walkthrough screenshots
    total = len(assets)
    for i, a in enumerate(assets, start=1):
        slide_screenshot(prs, i, total, a, captions.get(a.path.name, ""))

    # Metric 2 (out-of-the-box)
    slide_metric(prs, METRICS[1][0], METRICS[1][1], METRICS[1][2])

    # Metric 3 (internationalised)
    slide_metric(prs, METRICS[2][0], METRICS[2][1], METRICS[2][2])

    # Closers
    slide_cta(prs, video_url)
    slide_contact(prs)

    cp = prs.core_properties
    cp.author = "Idris-Michael Bakare"
    cp.title = "Agent-Reach -- Portfolio Deck"
    cp.modified = dt.datetime(2026, 5, 23, 0, 0, 0)
    cp.created  = dt.datetime(2026, 5, 23, 0, 0, 0)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"OK -> {OUTPUT}")
    print(f"   {len(prs.slides)} slides | {len(assets)} screenshots | captions: {CAPTIONS_FILE}")
    if not video_url:
        print("   video URL placeholder -> rerun with --video-url <url>")

def main() -> None:
    ap = argparse.ArgumentParser(description="Build the Agent-Reach showcase deck.")
    ap.add_argument("--video-url", default=None, help="YouTube/Loom URL for the CTA slide.")
    args = ap.parse_args()
    build(args.video_url)

if __name__ == "__main__":
    main()
