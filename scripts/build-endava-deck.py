"""
Build the AI-Native Engineering Portfolio deck (.pptx) for Idris-Michael Bakare.

Run:
    py -3.10 scripts/build-endava-deck.py
    py -3.10 scripts/build-endava-deck.py --video-url https://youtu.be/XXX

Inputs:
    Presentation slide deck assets - endava/*.png   (captured via Playwright)
    scripts/endava-captions.json                    (auto-seeded; user can edit)

Output:
    deliverables/endava-portfolio-deck.pptx
"""
from __future__ import annotations

import argparse
import datetime as dt
import json
import sys
from dataclasses import dataclass
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# Paths
ROOT = Path(__file__).resolve().parent.parent
ASSETS_DIR = ROOT / "Presentation slide deck assets - endava"
CAPTIONS_FILE = ROOT / "scripts" / "endava-captions.json"
OUTPUT = ROOT / "deliverables" / "endava-portfolio-deck.pptx"

# Design system - identical to Legal Pioneer deck
def hex_rgb(s: str) -> RGBColor:
    s = s.lstrip("#")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))

BG       = hex_rgb("#060610")
SURFACE  = hex_rgb("#0E0E1A")
ACCENT   = hex_rgb("#D4A017")
TEXT     = hex_rgb("#E2E8F0")
TEXT_DIM = hex_rgb("#94A3B8")
RULE     = hex_rgb("#1F1F2E")
GREEN    = hex_rgb("#4ADE80")
AMBER    = hex_rgb("#FBBF24")
RED      = hex_rgb("#F87171")
VIOLET   = hex_rgb("#C4B5FD")

FONT_DISPLAY = "Inter SemiBold"
FONT_BODY    = "Inter"
FONT_MONO    = "JetBrains Mono"

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# Project copy
PROJECT_TITLE     = "AI-Native Engineering Portfolio"
PROJECT_SUBTITLE  = "GA4 + Ads Modernisation Stack"
PROJECT_CONTEXT   = "Reference architecture for senior consulting engagements - Idris-Michael Bakare - Severus Connects, London"

THESIS = (
    "Senior consulting work fails when the spec is opinions\n"
    "and the artefact is a PDF."
)
THESIS_SUB = (
    "This portfolio is a working set of interactive dashboards that any prospective client "
    "can open in a browser and verify. System design, GA4 attribution, ads restructure, "
    "and stack mastery - delivered as auditable artefacts, not slideware."
)

METRICS = [
    ("8",       "Interactive HTML dashboards - opens in any browser, zero setup", "Dashboards shipped"),
    ("~150 KB", "Total payload across the entire portfolio - self-contained, no build step, no dependencies", "Total footprint"),
    ("5",       "Domains covered - system architecture, GA4 attribution, Google Ads restructure, tech stack mastery, client onboarding", "Domains covered"),
]

PILLARS = [
    ("Real artefacts, not slides",
     "Every claim in this portfolio resolves to an interactive dashboard the reviewer can open, "
     "inspect, and stress-test. No screenshots of work that doesn't exist."),
    ("Browser-native, zero setup",
     "Pure HTML, CSS, and vanilla JavaScript. No framework lock-in, no npm install, no build step. "
     "Double-click and it runs - on the reviewer's machine, today."),
    ("Auditable claims",
     "Every metric, every architectural decision, every campaign restructure is traceable. "
     "Hover any node, click any tier, read the rationale."),
    ("Reusable across engagements",
     "The same architectural patterns slot into Shopify, B2B SaaS, or services accounts. "
     "Designed once as a reference, deployed many times in client work."),
]

# Coverage map - replaces the Legal-Pioneer tier table
COVERAGE_MAP = [
    ("System Architecture", "End-to-end data flow from site -> GTM -> GA4 -> BigQuery -> Looker, with auth boundaries and failure modes called out.", "01"),
    ("Google Ads Case Study", "Performance Max + Search restructure: SKAG decomposition, asset group strategy, negative-keyword hygiene.", "02"),
    ("Ads Restructure", "Before/after account skeleton with budget reallocation rationale and bid strategy migration path.", "03"),
    ("GA4 Attribution", "Cross-channel attribution model with data-driven conversion paths, consent mode v2, and server-side measurement.", "04"),
    ("Tech Stack Mastery", "Stack inventory with depth-of-knowledge scoring across measurement, paid media, and analytics tooling.", "05"),
]

STACK_CHIPS = [
    "HTML5", "CSS3", "Vanilla JavaScript",
    "Google Analytics 4", "Google Tag Manager", "Google Ads",
    "Performance Max", "Data Layer", "Smart Bidding",
    "Consent Mode v2", "BigQuery", "Looker Studio",
]

CONTACT = {
    "name":     "Idris-Michael Bakare",
    "role":     "AI-Native Engineer - Severus Connects - London",
    "email":    "hello@severusconnects.com",
    "linkedin": "linkedin.com/in/idris-michael-bakare",
    "cta":      "Available to walk through the architecture live.",
}

# Helpers
@dataclass(frozen=True)
class Asset:
    path: Path
    width: int
    height: int

def discover_assets() -> list[Asset]:
    if not ASSETS_DIR.exists():
        sys.exit(f"Assets folder not found: {ASSETS_DIR}\nRun capture-endava-assets.py first.")
    pngs = sorted(ASSETS_DIR.glob("*.png"), key=lambda p: p.name)
    if not pngs:
        sys.exit(f"No PNGs found in: {ASSETS_DIR}")
    out: list[Asset] = []
    for p in pngs:
        with Image.open(p) as im:
            out.append(Asset(path=p, width=im.width, height=im.height))
    return out

DEFAULT_CAPTIONS: dict[str, str] = {
    "01-start-here.png":            "START_HERE - portfolio index. The single page a reviewer opens first to navigate every dashboard in the set.",
    "02-system-architecture.png":   "System Architecture - end-to-end data flow from site through GTM, GA4, BigQuery, and Looker. Auth boundaries and failure modes annotated inline.",
    "03-ads-dashboard.png":         "Ads Dashboard - account-level KPI roll-up across Search, Performance Max, and Shopping. Single view for weekly client reviews.",
    "04-google-ads-case-study.png": "Google Ads Case Study - SKAG decomposition into asset groups, negative-keyword hygiene, and bid-strategy migration. Before/after with rationale.",
    "05-ads-restructure.png":       "Ads Restructure - account skeleton redesign with budget reallocation and the audit log behind every decision.",
    "06-ga4-attribution.png":       "GA4 Attribution - data-driven model with consent mode v2, server-side measurement, and cross-channel conversion paths.",
    "07-tech-stack-mastery.png":    "Tech Stack Mastery - depth-of-knowledge scoring across measurement, paid media, and analytics. Honest, not aspirational.",
    "08-codecrafthub-dashboard.png":"CodeCraftHub Dashboard - client-facing onboarding view. Demonstrates the white-label surface clients see during a 90-day engagement.",
}

def load_or_seed_captions(assets: list[Asset]) -> dict[str, str]:
    if CAPTIONS_FILE.exists():
        try:
            existing: dict[str, str] = json.loads(CAPTIONS_FILE.read_text(encoding="utf-8"))
        except json.JSONDecodeError:
            existing = {}
    else:
        existing = {}
    seeded = {a.path.name: existing.get(a.path.name) or DEFAULT_CAPTIONS.get(a.path.name, "")
              for a in assets}
    CAPTIONS_FILE.parent.mkdir(parents=True, exist_ok=True)
    CAPTIONS_FILE.write_text(json.dumps(seeded, indent=2, sort_keys=True) + "\n", encoding="utf-8")
    return seeded

def fill_solid(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_text(slide, left_in: float, top_in: float, width_in: float, height_in: float,
             text: str, *, font: str = FONT_BODY, size_pt: int = 18, color: RGBColor = TEXT,
             bold: bool = False, align: int = PP_ALIGN.LEFT, anchor: int = MSO_ANCHOR.TOP) -> None:
    tb = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size_pt)
    f.bold = bold
    f.color.rgb = color

def paint_bg(slide) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SLIDE_W_IN), Inches(SLIDE_H_IN))
    fill_solid(bg, BG)

def add_accent_bar(slide, left_in: float = 0.7, top_in: float = 0.7, width_in: float = 0.6) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left_in), Inches(top_in), Inches(width_in), Inches(0.08))
    fill_solid(bar, ACCENT)

def add_page_footer(slide, label: str) -> None:
    add_text(slide, 0.5, SLIDE_H_IN - 0.45, SLIDE_W_IN - 1.0, 0.3,
             label, font=FONT_MONO, size_pt=9, color=TEXT_DIM)

# Slide builders
def slide_cover(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    add_accent_bar(s, 0.7, 2.2, 1.6)
    add_text(s, 0.7, 2.5, 12, 1.4, PROJECT_TITLE,
             font=FONT_DISPLAY, size_pt=64, color=TEXT, bold=True)
    add_text(s, 0.7, 3.95, 12, 0.7, PROJECT_SUBTITLE,
             font=FONT_DISPLAY, size_pt=26, color=ACCENT)
    add_text(s, 0.7, 4.8, 12, 0.8,
             "A reference architecture for senior consulting engagements -\n"
             "system design, attribution, ads restructure, and stack mastery.",
             font=FONT_BODY, size_pt=16, color=TEXT_DIM)
    add_text(s, 0.7, 6.6, 12, 0.4, PROJECT_CONTEXT,
             font=FONT_MONO, size_pt=11, color=TEXT_DIM)

def slide_thesis(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    add_accent_bar(s)
    add_text(s, 0.7, 1.0, 12, 0.6, "01 - The problem",
             font=FONT_MONO, size_pt=12, color=ACCENT)
    add_text(s, 0.7, 2.0, 12, 3.0, THESIS,
             font=FONT_DISPLAY, size_pt=40, color=TEXT, bold=True)
    add_text(s, 0.7, 5.4, 12, 1.6, THESIS_SUB,
             font=FONT_BODY, size_pt=18, color=TEXT_DIM)
    add_page_footer(s, "portfolio - thesis")

def slide_pillars(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    add_accent_bar(s)
    add_text(s, 0.7, 1.0, 12, 0.6, "02 - Design principles",
             font=FONT_MONO, size_pt=12, color=ACCENT)
    add_text(s, 0.7, 1.5, 12, 0.7, "Four design decisions.",
             font=FONT_DISPLAY, size_pt=32, color=TEXT, bold=True)
    coords = [(0.7, 2.6), (6.9, 2.6), (0.7, 4.95), (6.9, 4.95)]
    for (title, body), (x, y) in zip(PILLARS, coords):
        card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(5.7), Inches(2.1))
        card.fill.solid()
        card.fill.fore_color.rgb = SURFACE
        card.line.color.rgb = RULE
        card.line.width = Pt(0.75)
        add_text(s, x + 0.3, y + 0.25, 5.3, 0.5, title,
                 font=FONT_DISPLAY, size_pt=18, color=ACCENT, bold=True)
        add_text(s, x + 0.3, y + 0.85, 5.3, 1.2, body,
                 font=FONT_BODY, size_pt=13, color=TEXT)
    add_page_footer(s, "portfolio - principles")

def slide_coverage_map(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    add_accent_bar(s)
    add_text(s, 0.7, 1.0, 12, 0.6, "03 - Coverage map",
             font=FONT_MONO, size_pt=12, color=ACCENT)
    add_text(s, 0.7, 1.5, 12, 0.7, "Five domains. Eight dashboards.",
             font=FONT_DISPLAY, size_pt=28, color=TEXT, bold=True)

    headers = ["#", "Domain", "What it demonstrates"]
    col_x = [0.7, 1.6, 5.6]
    col_w = [0.7, 3.8, 7.0]
    header_top = 3.0
    for h, x, w in zip(headers, col_x, col_w):
        add_text(s, x, header_top, w, 0.4, h, font=FONT_MONO, size_pt=11, color=TEXT_DIM)
    rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(header_top + 0.45),
                               Inches(SLIDE_W_IN - 1.4), Inches(0.02))
    fill_solid(rule, RULE)

    row_h = 0.75
    for i, (domain, demo, num) in enumerate(COVERAGE_MAP):
        y = header_top + 0.7 + i * row_h
        add_text(s, col_x[0], y + 0.05, col_w[0], 0.5, num,
                 font=FONT_MONO, size_pt=14, color=ACCENT, bold=True)
        add_text(s, col_x[1], y + 0.05, col_w[1], 0.5, domain,
                 font=FONT_DISPLAY, size_pt=15, color=TEXT, bold=True)
        add_text(s, col_x[2], y + 0.05, col_w[2], 0.7, demo,
                 font=FONT_BODY, size_pt=12, color=TEXT_DIM)
    add_page_footer(s, "portfolio - coverage")

def slide_metric(prs: Presentation, big: str, sub: str, label: str) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    add_accent_bar(s)
    add_text(s, 0.7, 1.0, 12, 0.6, label,
             font=FONT_MONO, size_pt=12, color=ACCENT)
    # Smaller font for longer metric strings
    big_size = 130 if len(big) <= 4 else 96
    add_text(s, 0.7, 2.4, 12, 2.5, big,
             font=FONT_DISPLAY, size_pt=big_size, color=ACCENT, bold=True)
    add_text(s, 0.7, 5.5, 12, 1.5, sub,
             font=FONT_BODY, size_pt=20, color=TEXT)
    add_page_footer(s, "portfolio - impact")

def slide_screenshot(prs: Presentation, index: int, total: int, asset: Asset, caption: str) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    img_area_top = 0.4
    img_area_height = SLIDE_H_IN - img_area_top - 0.9
    img_area_width = SLIDE_W_IN - 1.0
    img_area_left = 0.5

    aspect = asset.width / asset.height
    area_aspect = img_area_width / img_area_height
    if aspect >= area_aspect:
        new_w = img_area_width
        new_h = new_w / aspect
    else:
        new_h = img_area_height
        new_w = new_h * aspect
    left = img_area_left + (img_area_width - new_w) / 2
    top = img_area_top + (img_area_height - new_h) / 2
    s.shapes.add_picture(str(asset.path), Inches(left), Inches(top), Inches(new_w), Inches(new_h))

    strip_top = SLIDE_H_IN - 0.85
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.5), Inches(strip_top),
                              Inches(0.08), Inches(0.45))
    fill_solid(bar, ACCENT)
    add_text(s, 0.75, strip_top - 0.05, 11.5, 0.55,
             caption if caption else "-",
             font=FONT_DISPLAY, size_pt=14,
             color=TEXT if caption else TEXT_DIM)
    add_text(s, SLIDE_W_IN - 1.2, strip_top + 0.05, 0.8, 0.3,
             f"{index:02d}/{total:02d}",
             font=FONT_MONO, size_pt=10, color=TEXT_DIM, align=PP_ALIGN.RIGHT)

def slide_tech_chips(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    add_accent_bar(s)
    add_text(s, 0.7, 1.0, 12, 0.6, "Stack proof",
             font=FONT_MONO, size_pt=12, color=ACCENT)
    add_text(s, 0.7, 1.5, 12, 0.8, "What it's built with.",
             font=FONT_DISPLAY, size_pt=36, color=TEXT, bold=True)

    chip_h = 0.5
    pad_x = 0.2
    pad_y = 0.2
    char_w = 0.105
    left = 0.7
    top = 2.8
    cur_x = left
    cur_y = top
    max_x = SLIDE_W_IN - 0.7

    for chip in STACK_CHIPS:
        w = max(1.1, len(chip) * char_w + 0.5)
        if cur_x + w > max_x:
            cur_x = left
            cur_y += chip_h + pad_y
        rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(cur_x), Inches(cur_y), Inches(w), Inches(chip_h))
        rect.fill.solid()
        rect.fill.fore_color.rgb = SURFACE
        rect.line.color.rgb = ACCENT
        rect.line.width = Pt(0.75)
        tf = rect.text_frame
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = chip
        run.font.name = FONT_DISPLAY
        run.font.size = Pt(13)
        run.font.color.rgb = ACCENT
        cur_x += w + pad_x
    add_page_footer(s, "portfolio - stack")

def slide_cta(prs: Presentation, video_url: str | None) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    add_accent_bar(s)
    add_text(s, 0.7, 1.0, 12, 0.6, "Watch the walkthrough",
             font=FONT_MONO, size_pt=12, color=ACCENT)
    add_text(s, 0.7, 2.0, 12, 1.5, "See the portfolio live.",
             font=FONT_DISPLAY, size_pt=52, color=TEXT, bold=True)
    add_text(s, 0.7, 3.6, 12, 1.2,
             "Eight dashboards. One browser tab. Every architectural\n"
             "decision visible, every claim auditable.",
             font=FONT_BODY, size_pt=18, color=TEXT_DIM)
    url = video_url or "YOUTUBE_URL_PENDING"
    add_text(s, 0.7, 5.3, 12, 0.5, "Link:",
             font=FONT_MONO, size_pt=12, color=TEXT_DIM)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(5.7), Inches(12), Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = url
    run.font.name = FONT_MONO
    run.font.size = Pt(22)
    run.font.color.rgb = ACCENT
    run.font.bold = True
    if video_url:
        run.hyperlink.address = video_url
    add_page_footer(s, "portfolio - cta")

def slide_contact(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    add_accent_bar(s)
    add_text(s, 0.7, 1.0, 12, 0.6, "Contact",
             font=FONT_MONO, size_pt=12, color=ACCENT)
    add_text(s, 0.7, 2.2, 12, 1.0, CONTACT["name"],
             font=FONT_DISPLAY, size_pt=54, color=TEXT, bold=True)
    add_text(s, 0.7, 3.4, 12, 0.7, CONTACT["role"],
             font=FONT_BODY, size_pt=22, color=TEXT_DIM)
    add_text(s, 0.7, 5.0, 12, 0.5, CONTACT["email"],
             font=FONT_MONO, size_pt=18, color=ACCENT)
    add_text(s, 0.7, 5.6, 12, 0.5, CONTACT["linkedin"],
             font=FONT_MONO, size_pt=18, color=ACCENT)
    add_text(s, 0.7, 6.6, 12, 0.4, CONTACT["cta"],
             font=FONT_DISPLAY, size_pt=18, color=TEXT)

# Build
def build(video_url: str | None) -> None:
    assets = discover_assets()
    captions = load_or_seed_captions(assets)

    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    # Framing
    slide_cover(prs)
    slide_thesis(prs)
    slide_pillars(prs)
    slide_coverage_map(prs)

    # Metric 1 (dashboards)
    slide_metric(prs, METRICS[0][0], METRICS[0][1], METRICS[0][2])

    # Visual proof - dashboard walkthrough
    total = len(assets)
    for i, a in enumerate(assets, start=1):
        slide_screenshot(prs, i, total, a, captions.get(a.path.name, ""))

    # Metric 2 (footprint)
    slide_metric(prs, METRICS[1][0], METRICS[1][1], METRICS[1][2])

    # Metric 3 (domains)
    slide_metric(prs, METRICS[2][0], METRICS[2][1], METRICS[2][2])

    # Closers
    slide_tech_chips(prs)
    slide_cta(prs, video_url)
    slide_contact(prs)

    cp = prs.core_properties
    cp.author = "Idris-Michael Bakare"
    cp.title = "AI-Native Engineering Portfolio - GA4 + Ads Modernisation Stack"
    cp.modified = dt.datetime(2026, 5, 23, 0, 0, 0)
    cp.created  = dt.datetime(2026, 5, 23, 0, 0, 0)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"OK -> {OUTPUT}")
    print(f"   {len(prs.slides)} slides | {len(assets)} screenshots | captions: {CAPTIONS_FILE}")
    if not video_url:
        print("   video URL placeholder -> rerun with --video-url <url>")

def main() -> None:
    ap = argparse.ArgumentParser(description="Build the AI-Native Engineering Portfolio deck.")
    ap.add_argument("--video-url", default=None, help="YouTube/Loom URL for the CTA slide.")
    args = ap.parse_args()
    build(args.video_url)

if __name__ == "__main__":
    main()
