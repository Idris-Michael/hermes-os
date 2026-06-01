"""
Build the Hermes OS LinkedIn showcase deck (.pptx).

Run:
    py -3.10 scripts/build-linkedin-deck.py
    py -3.10 scripts/build-linkedin-deck.py --video-url https://youtu.be/XXX

Inputs:
    Presentation slide deck assets/*.png  (31 screenshots, one slide each)
    severus-connects-prompts/linkedin-projects-optimised.md  (copy source, not parsed — used as design reference)
    scripts/deck-captions.json  (auto-generated on first run; user edits to add captions)

Output:
    deliverables/hermes-os-linkedin-deck.pptx

Deterministic: same inputs + same caption file => byte-identical .pptx (modulo pptx's internal timestamps,
which are zeroed via core_properties.modified override).
"""
from __future__ import annotations

import argparse
import datetime as dt
import json
import sys
from dataclasses import dataclass
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# ─── Paths ─────────────────────────────────────────────────────────────────────
ROOT = Path(__file__).resolve().parent.parent
ASSETS_DIR = ROOT / "Presentation slide deck assets"
CAPTIONS_FILE = ROOT / "scripts" / "deck-captions.json"
OUTPUT = ROOT / "deliverables" / "hermes-os-linkedin-deck.pptx"

# ─── Design system (single source of truth) ────────────────────────────────────
def hex_rgb(s: str) -> RGBColor:
    s = s.lstrip("#")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))

BG        = hex_rgb("#060610")
SURFACE   = hex_rgb("#0E0E1A")
ACCENT    = hex_rgb("#D4A017")
TEXT      = hex_rgb("#E2E8F0")
TEXT_DIM  = hex_rgb("#94A3B8")
RULE      = hex_rgb("#1F1F2E")

FONT_DISPLAY = "Inter SemiBold"
FONT_BODY    = "Inter"
FONT_MONO    = "JetBrains Mono"

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# ─── LinkedIn project copy (mirrors linkedin-projects-optimised.md) ────────────
PROJECT_METRICS = [
    ("$0.25 / month",      "Operating cost of 5 coordinated services on one machine",      "Hermes Hub — Agency Automation OS"),
    ("1.4× → 3.8× ROAS",   "Shopify brand, 8 weeks, anti-gravity quality-gate framework",  "GA4 + Google Ads Delivery System"),
    ("~$0.21 / month",     "Fully automated Instagram pipeline for two accounts",          "Instagram AI Content Pipeline"),
]

TECH_CHIPS = [
    "React 19", "TypeScript", "Express.js", "SQLite", "Vite",
    "Three.js", "React Three Fiber", "Electron", "Tailwind CSS",
    "Hugging Face Inference", "Qwen3-32B", "DeepSeek V3",
    "Google Gemini TTS", "Hyperframes", "Meta Graph API",
    "Telegram Bot API", "Google Analytics 4", "Google Ads",
    "GTM", "Performance Max", "Smart Bidding", "Merchant Centre",
    "Cron Scheduling", "REST API Design", "Monorepo Architecture",
    "Multi-Agent Orchestration",
]

CONTACT = {
    "name":     "Idris-Michael Bakare",
    "agency":   "Severus Connects · London",
    "email":    "hello@severusconnects.com",
    "linkedin": "linkedin.com/in/idris-michael-bakare",
    "cta":      "DM to scope a project",
}

# ─── Helpers ───────────────────────────────────────────────────────────────────
@dataclass(frozen=True)
class Asset:
    path: Path
    width: int
    height: int

def discover_assets() -> list[Asset]:
    if not ASSETS_DIR.exists():
        sys.exit(f"Assets folder not found: {ASSETS_DIR}")
    pngs = sorted(ASSETS_DIR.glob("*.png"), key=lambda p: p.name)
    if not pngs:
        sys.exit(f"No PNGs found in: {ASSETS_DIR}")
    out: list[Asset] = []
    for p in pngs:
        with Image.open(p) as im:
            out.append(Asset(path=p, width=im.width, height=im.height))
    return out

def load_or_seed_captions(assets: list[Asset]) -> dict[str, str]:
    if CAPTIONS_FILE.exists():
        try:
            existing: dict[str, str] = json.loads(CAPTIONS_FILE.read_text(encoding="utf-8"))
        except json.JSONDecodeError:
            existing = {}
    else:
        existing = {}
    # Add any new files with empty caption; preserve user edits.
    seeded = {a.path.name: existing.get(a.path.name, "") for a in assets}
    CAPTIONS_FILE.parent.mkdir(parents=True, exist_ok=True)
    CAPTIONS_FILE.write_text(json.dumps(seeded, indent=2, sort_keys=True) + "\n", encoding="utf-8")
    return seeded

def fill_solid(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_rect(slide, left_in: float, top_in: float, width_in: float, height_in: float, color: RGBColor):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
    fill_solid(rect, color)
    return rect

def add_text(slide, left_in: float, top_in: float, width_in: float, height_in: float,
             text: str, *, font: str = FONT_BODY, size_pt: int = 18, color: RGBColor = TEXT,
             bold: bool = False, align: int = PP_ALIGN.LEFT, anchor: int = MSO_ANCHOR.TOP) -> None:
    tb = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size_pt)
    f.bold = bold
    f.color.rgb = color

def paint_bg(slide) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SLIDE_W_IN), Inches(SLIDE_H_IN))
    fill_solid(bg, BG)

def add_accent_bar(slide, left_in: float = 0.5, top_in: float = 0.5, width_in: float = 0.6) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left_in), Inches(top_in), Inches(width_in), Inches(0.08))
    fill_solid(bar, ACCENT)

def add_page_footer(slide, label: str) -> None:
    add_text(slide, 0.5, SLIDE_H_IN - 0.45, SLIDE_W_IN - 1.0, 0.3,
             label, font=FONT_MONO, size_pt=9, color=TEXT_DIM)

# ─── Slide builders ────────────────────────────────────────────────────────────
def slide_cover(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    add_accent_bar(s, 0.7, 2.2, 1.6)
    add_text(s, 0.7, 2.5, 12, 1.4,
             "Hermes OS",
             font=FONT_DISPLAY, size_pt=80, color=TEXT, bold=True)
    add_text(s, 0.7, 3.8, 12, 0.7,
             "An Agency Automation Operating System",
             font=FONT_DISPLAY, size_pt=28, color=ACCENT)
    add_text(s, 0.7, 4.7, 12, 0.5,
             "Built and operated solo · London · 2026",
             font=FONT_BODY, size_pt=16, color=TEXT_DIM)
    add_text(s, 0.7, 6.6, 12, 0.4,
             "Idris-Michael Bakare · Severus Connects",
             font=FONT_MONO, size_pt=11, color=TEXT_DIM)

def slide_thesis(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    add_accent_bar(s, 0.7, 0.7)
    add_text(s, 0.7, 1.0, 12, 0.6, "01 · Thesis",
             font=FONT_MONO, size_pt=12, color=ACCENT)
    add_text(s, 0.7, 2.2, 12, 3.5,
             "One-person agencies don't lose to capacity.\nThey lose to context-switching.",
             font=FONT_DISPLAY, size_pt=48, color=TEXT, bold=True)
    add_text(s, 0.7, 5.4, 12, 1.2,
             "Hermes OS is the infrastructure I built so the agency runs on rails — "
             "intake, proposal, delivery, reporting, and approval all route through one coherent system.",
             font=FONT_BODY, size_pt=18, color=TEXT_DIM)
    add_page_footer(s, "hermes-os · thesis")

def slide_agenda(prs: Presentation, screenshot_count: int) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    add_accent_bar(s, 0.7, 0.7)
    add_text(s, 0.7, 1.0, 12, 0.6, "02 · What you're about to see",
             font=FONT_MONO, size_pt=12, color=ACCENT)
    items = [
        "Stack overview & coordinated services",
        f"Hermes Hub walkthrough — {screenshot_count} live UI captures",
        "Project metrics (operating cost · ROAS · automation reach)",
        "Tech stack proof",
        "Watch the full video walkthrough",
    ]
    for i, item in enumerate(items):
        top = 2.3 + i * 0.85
        add_text(s, 0.9, top, 0.5, 0.6, f"0{i+1}",
                 font=FONT_MONO, size_pt=20, color=ACCENT, bold=True)
        add_text(s, 1.7, top + 0.05, 11, 0.6, item,
                 font=FONT_DISPLAY, size_pt=22, color=TEXT)
    add_page_footer(s, "hermes-os · agenda")

def slide_stack(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    add_accent_bar(s, 0.7, 0.7)
    add_text(s, 0.7, 1.0, 12, 0.6, "03 · Stack",
             font=FONT_MONO, size_pt=12, color=ACCENT)
    add_text(s, 0.7, 1.5, 12, 0.8, "Five services. One machine.",
             font=FONT_DISPLAY, size_pt=36, color=TEXT, bold=True)
    services = [
        ("Hermes Hub",        "React 19 · TypeScript · Express · SQLite — command layer"),
        ("severus-social",    "Instagram AI pipeline · Qwen3 · Gemini TTS · Hyperframes"),
        ("Telegram Gateway",  "Human-in-the-loop approval for every publishable artefact"),
        ("UGC Generator",     "Local Electron app for video variants and brand assets"),
        ("Knowledge Base",    "Three.js / R3F 3D constellation of project context"),
    ]
    for i, (name, desc) in enumerate(services):
        top = 2.7 + i * 0.75
        add_text(s, 0.7, top, 3.5, 0.5, name,
                 font=FONT_DISPLAY, size_pt=20, color=ACCENT, bold=True)
        add_text(s, 4.4, top + 0.05, 8.5, 0.5, desc,
                 font=FONT_BODY, size_pt=16, color=TEXT)
    add_page_footer(s, "hermes-os · stack")

def slide_screenshot(prs: Presentation, index: int, total: int, asset: Asset, caption: str) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    # Reserve a caption strip at the bottom (~0.9in). Image fills the rest with letterboxing.
    img_area_top = 0.4
    img_area_height = SLIDE_H_IN - img_area_top - 0.9
    img_area_width = SLIDE_W_IN - 1.0
    img_area_left = 0.5

    # Fit image preserving aspect ratio (letterbox on the BG).
    aspect = asset.width / asset.height
    area_aspect = img_area_width / img_area_height
    if aspect >= area_aspect:
        # Image is wider — fit by width.
        new_w = img_area_width
        new_h = new_w / aspect
    else:
        new_h = img_area_height
        new_w = new_h * aspect
    left = img_area_left + (img_area_width - new_w) / 2
    top = img_area_top + (img_area_height - new_h) / 2
    s.shapes.add_picture(str(asset.path), Inches(left), Inches(top), Inches(new_w), Inches(new_h))

    # Caption strip
    strip_top = SLIDE_H_IN - 0.85
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.5), Inches(strip_top),
                              Inches(0.08), Inches(0.45))
    fill_solid(bar, ACCENT)
    add_text(s, 0.75, strip_top - 0.05, 11.5, 0.55,
             caption if caption else "—",
             font=FONT_DISPLAY, size_pt=14,
             color=TEXT if caption else TEXT_DIM)
    add_text(s, SLIDE_W_IN - 1.2, strip_top + 0.05, 0.8, 0.3,
             f"{index:02d}/{total:02d}",
             font=FONT_MONO, size_pt=10, color=TEXT_DIM, align=PP_ALIGN.RIGHT)

def slide_metric(prs: Presentation, metric: str, sub: str, project_title: str, n: int) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    add_accent_bar(s, 0.7, 0.7)
    add_text(s, 0.7, 1.0, 12, 0.6, f"Project 0{n} · {project_title}",
             font=FONT_MONO, size_pt=12, color=ACCENT)
    add_text(s, 0.7, 2.5, 12, 2.5, metric,
             font=FONT_DISPLAY, size_pt=110, color=ACCENT, bold=True)
    add_text(s, 0.7, 5.5, 12, 1.2, sub,
             font=FONT_BODY, size_pt=22, color=TEXT)
    add_page_footer(s, f"project · 0{n}")

def slide_tech_chips(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    add_accent_bar(s, 0.7, 0.7)
    add_text(s, 0.7, 1.0, 12, 0.6, "Stack proof",
             font=FONT_MONO, size_pt=12, color=ACCENT)
    add_text(s, 0.7, 1.5, 12, 0.8, "What it's built with.",
             font=FONT_DISPLAY, size_pt=36, color=TEXT, bold=True)

    # Layout chips in rows, wrap as needed.
    chip_h = 0.5
    pad_x = 0.2
    pad_y = 0.2
    char_w = 0.105  # rough Inter SemiBold 14pt char width in inches
    left = 0.7
    top = 2.8
    cur_x = left
    cur_y = top
    max_x = SLIDE_W_IN - 0.7

    for chip in TECH_CHIPS:
        w = max(1.1, len(chip) * char_w + 0.5)
        if cur_x + w > max_x:
            cur_x = left
            cur_y += chip_h + pad_y
        rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(cur_x), Inches(cur_y), Inches(w), Inches(chip_h))
        rect.fill.solid()
        rect.fill.fore_color.rgb = SURFACE
        rect.line.color.rgb = ACCENT
        rect.line.width = Pt(0.75)
        tf = rect.text_frame
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = chip
        run.font.name = FONT_DISPLAY
        run.font.size = Pt(13)
        run.font.color.rgb = ACCENT
        cur_x += w + pad_x

    add_page_footer(s, "hermes-os · stack")

def slide_cta(prs: Presentation, video_url: str | None) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    add_accent_bar(s, 0.7, 0.7)
    add_text(s, 0.7, 1.0, 12, 0.6, "Watch the full walkthrough",
             font=FONT_MONO, size_pt=12, color=ACCENT)
    add_text(s, 0.7, 2.0, 12, 1.5,
             "See it running, end-to-end.",
             font=FONT_DISPLAY, size_pt=52, color=TEXT, bold=True)
    add_text(s, 0.7, 3.6, 12, 1,
             "A 3-minute walkthrough of the dashboard, the agent system,\n"
             "and the Instagram pipeline producing a real post.",
             font=FONT_BODY, size_pt=18, color=TEXT_DIM)

    url = video_url or "YOUTUBE_URL_PENDING"
    add_text(s, 0.7, 5.3, 12, 0.5, "Link:",
             font=FONT_MONO, size_pt=12, color=TEXT_DIM)
    # Insert as hyperlink so it's clickable in the exported PDF.
    tb = s.shapes.add_textbox(Inches(0.7), Inches(5.7), Inches(12), Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = url
    run.font.name = FONT_MONO
    run.font.size = Pt(22)
    run.font.color.rgb = ACCENT
    run.font.bold = True
    if video_url:
        run.hyperlink.address = video_url

    add_page_footer(s, "hermes-os · cta")

def slide_contact(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s)
    add_accent_bar(s, 0.7, 0.7)
    add_text(s, 0.7, 1.0, 12, 0.6, "Contact",
             font=FONT_MONO, size_pt=12, color=ACCENT)
    add_text(s, 0.7, 2.2, 12, 1.0, CONTACT["name"],
             font=FONT_DISPLAY, size_pt=54, color=TEXT, bold=True)
    add_text(s, 0.7, 3.4, 12, 0.7, CONTACT["agency"],
             font=FONT_BODY, size_pt=22, color=TEXT_DIM)
    add_text(s, 0.7, 5.0, 12, 0.5, CONTACT["email"],
             font=FONT_MONO, size_pt=18, color=ACCENT)
    add_text(s, 0.7, 5.6, 12, 0.5, CONTACT["linkedin"],
             font=FONT_MONO, size_pt=18, color=ACCENT)
    add_text(s, 0.7, 6.6, 12, 0.4, CONTACT["cta"],
             font=FONT_DISPLAY, size_pt=18, color=TEXT)

# ─── Build ─────────────────────────────────────────────────────────────────────
def build(video_url: str | None) -> None:
    assets = discover_assets()
    captions = load_or_seed_captions(assets)

    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    # Framing
    slide_cover(prs)
    slide_thesis(prs)
    slide_agenda(prs, len(assets))
    slide_stack(prs)

    # Project 1 metric drop, then walkthrough screenshots
    slide_metric(prs, PROJECT_METRICS[0][0], PROJECT_METRICS[0][1], PROJECT_METRICS[0][2], 1)

    total = len(assets)
    for i, a in enumerate(assets, start=1):
        slide_screenshot(prs, i, total, a, captions.get(a.path.name, ""))

    # Project 2 & 3 metric drops (no native screenshots for these)
    slide_metric(prs, PROJECT_METRICS[1][0], PROJECT_METRICS[1][1], PROJECT_METRICS[1][2], 2)
    slide_metric(prs, PROJECT_METRICS[2][0], PROJECT_METRICS[2][1], PROJECT_METRICS[2][2], 3)

    # Closers
    slide_tech_chips(prs)
    slide_cta(prs, video_url)
    slide_contact(prs)

    # Zero out timestamps for deterministic output
    cp = prs.core_properties
    cp.author = "Severus Connects"
    cp.title = "Hermes OS — LinkedIn Showcase"
    cp.modified = dt.datetime(2026, 1, 1, 0, 0, 0)
    cp.created  = dt.datetime(2026, 1, 1, 0, 0, 0)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"OK · wrote {OUTPUT}")
    print(f"   · {len(prs.slides)} slides · {len(assets)} screenshots · captions file: {CAPTIONS_FILE}")
    if not video_url:
        print("   · video URL placeholder — rerun with --video-url <url> when video is uploaded")

def main() -> None:
    ap = argparse.ArgumentParser(description="Build the Hermes OS LinkedIn showcase deck.")
    ap.add_argument("--video-url", default=None, help="YouTube/Loom URL for the CTA slide.")
    args = ap.parse_args()
    build(args.video_url)

if __name__ == "__main__":
    main()
